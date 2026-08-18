import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pptx_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Dark Theme / Cyberpunk Cyan)
    BG_DARK = RGBColor(10, 15, 30)
    ACCENT_CYAN = RGBColor(0, 242, 254)
    ACCENT_PURPLE = RGBColor(127, 0, 255)
    TEXT_LIGHT = RGBColor(240, 246, 252)
    TEXT_MUTED = RGBColor(139, 148, 158)
    COLOR_SUCCESS = RGBColor(16, 185, 129)

    blank_layout = prs.slide_layouts[6]

    def add_dark_slide(title_text=""):
        slide = prs.slides.add_slide(blank_layout)
        
        # Dark Background shape
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()

        if title_text:
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.9))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = 'Arial'
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = ACCENT_CYAN

        return slide

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    slide1 = add_dark_slide()
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "AI-BASED VOICE ASSISTANT (ARIA)"
    p1.font.name = 'Arial'
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf1.add_paragraph()
    p2.text = "Speech Recognition, Google OAuth Sign-In, Email/Password Auth & SQLite Database"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_LIGHT
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf1.add_paragraph()
    p3.text = "\nMajor Academic Project Presentation  |  Python 3.13 • SQLite3 • Google GIS • Flask"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 2: Problem Statement & Key Objectives
    # ----------------------------------------------------
    slide2 = add_dark_slide("Problem Statement & Objectives")
    tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    items2 = [
        ("Problem Statement", "Modern AI assistants require hands-free interaction, secure multi-user authentication, and persistent data storage. ARIA solves this by providing voice synthesis, desktop automation, Google OAuth, Email/Password auth, and an SQLite database."),
        ("Key Objectives", "1. Google OAuth Sign-In & Email/Password SQL authentication.\n2. Persistent user account & chat history storage in SQLite database.\n3. Real-time microphone capture & speech-to-text conversion.\n4. Offline voice response synthesis via pyttsx3 & browser TTS.\n5. Automated system app launcher, math evaluator, and web search reports.")
    ]
    for heading, text in items2:
        p_head = tf2.add_paragraph()
        p_head.text = heading
        p_head.font.size = Pt(20)
        p_head.font.bold = True
        p_head.font.color.rgb = ACCENT_PURPLE

        p_desc = tf2.add_paragraph()
        p_desc.text = text
        p_desc.font.size = Pt(15)
        p_desc.font.color.rgb = TEXT_LIGHT
        p_desc.space_after = Pt(14)

    # ----------------------------------------------------
    # SLIDE 3: Newest Features (Google Auth & SQL DB)
    # ----------------------------------------------------
    slide3 = add_dark_slide("Newest Added Features (Google Auth & SQL DB)")
    tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf3 = tb3.text_frame
    tf3.word_wrap = True

    new_feats = [
        ("Google Sign-In OAuth System", "Official Google Identity Services (GIS SDK) integration. Supports 1-click Google Sign-In with JWT token decoding and profile picture sync."),
        ("Email & Password Auth System", "Complete registration (/api/auth/register) and login (/api/auth/email-login). Password credentials protected with salted Werkzeug cryptographic hashes."),
        ("Persistent Relational SQL Database", "SQLite database (database.db) with structured tables: 'users' (accounts & auth) and 'chat_history' (voice logs & chat sessions)."),
        ("Auth Mode Switcher UI", "Interactive tabs to switch between 'Sign In' and 'Create Account' with instant status notifications and guest fallback.")
    ]
    for title, desc in new_feats:
        p = tf3.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {title}: "
        r1.font.bold = True
        r1.font.size = Pt(17)
        r1.font.color.rgb = ACCENT_CYAN

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 4: SQL Database Schema Architecture
    # ----------------------------------------------------
    slide4 = add_dark_slide("SQL Relational Database Schema (database.db)")
    tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf4 = tb4.text_frame
    tf4.word_wrap = True

    db_tables = [
        ("Table 1: 'users' (User Accounts & Auth)", "Columns: id (PK), name, email (UNIQUE), password_hash (Werkzeug salted hash), provider ('Google'/'Email'), picture, created_at, last_login.\nStores registered user credentials securely."),
        ("Table 2: 'chat_history' (Conversation Logs)", "Columns: id (PK AUTO), user_id (FK), session_id, role ('user'/'assistant'), message, category, timestamp.\nPersists user voice transcripts and assistant responses linked to individual account IDs.")
    ]
    for t_name, t_desc in db_tables:
        p_head = tf4.add_paragraph()
        p_head.text = t_name
        p_head.font.size = Pt(19)
        p_head.font.bold = True
        p_head.font.color.rgb = ACCENT_PURPLE

        p_desc = tf4.add_paragraph()
        p_desc.text = t_desc
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = TEXT_LIGHT
        p_desc.space_after = Pt(16)

    # ----------------------------------------------------
    # SLIDE 5: Terms and Conditions
    # ----------------------------------------------------
    slide5 = add_dark_slide("Terms and Conditions (Operational & Privacy)")
    tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf5 = tb5.text_frame
    tf5.word_wrap = True

    terms = [
        ("User Privacy & Password Hashing", "Passwords are strictly stored as one-way PBKDF2/SHA-256 hashes. Google OAuth tokens are processed solely for identification and account linkage."),
        ("Microphone Stream Usage", "Microphone audio is transcribed locally or via browser Web Speech API solely for command execution. Raw audio streams are never stored permanently without consent."),
        ("Account Confidentiality", "Users are responsible for maintaining credential secrecy. Chat history logs are strictly linked to user IDs in the SQLite database to preserve session privacy."),
        ("Service Availability & Rate Limits", "External API services (Wikipedia API, Google Search) are accessed within rate limits to ensure robust uptime and prevent service degradation.")
    ]
    for term_title, term_desc in terms:
        p = tf5.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {term_title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = ACCENT_CYAN

        r2 = p.add_run()
        r2.text = term_desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 6: Error Rectification & Diagnostics
    # ----------------------------------------------------
    slide6 = add_dark_slide("Error Rectification & Diagnostics")
    tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf6 = tb6.text_frame
    tf6.word_wrap = True

    errors = [
        ("Database Schema Migration Rectification", "Added automatic PRAGMA table_info inspections and ALTER TABLE migrations in app.py to prevent 'missing column' crashes on legacy DB files."),
        ("Credential Validation Rectification", "Enforced password length constraints, duplicate email checks, and 401 Unauthorized responses to prevent corrupt database entries."),
        ("Google OAuth Popup Block Fallback", "Handled GIS SDK popup blockages with an interactive prompt dialog fallback so users can always sign in smoothly."),
        ("Web Speech Permission Error Handling", "Handled 'not-allowed' microphone errors by seamlessly falling back to keyboard command input with clear visual UI alerts.")
    ]
    for err_title, err_desc in errors:
        p = tf6.add_paragraph()
        r1 = p.add_run()
        r1.text = f"[FIXED] {err_title}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = COLOR_SUCCESS

        r2 = p.add_run()
        r2.text = err_desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 7: System Architecture & Tech Stack
    # ----------------------------------------------------
    slide7 = add_dark_slide("Technology Stack")
    tb7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf7 = tb7.text_frame
    tf7.word_wrap = True

    tech_stack = [
        ("Python 3.13", "Core backend language for assistant engine, intent routing, and API endpoints."),
        ("SQLite 3 Database", "Persistent SQL database engine storing user credentials and conversation sessions."),
        ("Google Identity Services SDK", "Official Google OAuth 2.0 single sign-on authentication SDK."),
        ("SpeechRecognition & pyttsx3", "Captures microphone speech and synthesizes offline spoken audio feedback."),
        ("Flask REST Framework", "Serves Web Control Center REST API at http://localhost:5000."),
        ("HTML5, CSS3, JavaScript", "Futuristic Glassmorphic Dark UI, Web Audio API sound effects, and interactive modals.")
    ]
    for tech, desc in tech_stack:
        p = tf7.add_paragraph()
        r1 = p.add_run()
        r1.text = f"- {tech}: "
        r1.font.bold = True
        r1.font.size = Pt(16)
        r1.font.color.rgb = ACCENT_CYAN

        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 8: Testing & Verification Results
    # ----------------------------------------------------
    slide8 = add_dark_slide("Testing & Verification (14/14 Tests Passed)")
    tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    tf8 = tb8.text_frame
    tf8.word_wrap = True

    p_t = tf8.add_paragraph()
    p_t.text = "Automated Unit Test Suites (test_assistant.py & test_auth_db.py):"
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = ACCENT_CYAN
    p_t.space_after = Pt(8)

    tests_list = [
        "[PASS] Voice Engine Tests 1-10: Time, Greeting, Stats, Math, Joke, Apps, Wiki -> 100% PASSED",
        "[PASS] SQL Test 1: SQLite Schema Integrity & Migration -> 100% PASSED",
        "[PASS] SQL Test 2: Email & Password User Account Registration -> 100% PASSED",
        "[PASS] SQL Test 3: Password Hashing & Email Authentication -> 100% PASSED",
        "[PASS] SQL Test 4: Google OAuth User Linking & DB Sync -> 100% PASSED"
    ]
    for t_item in tests_list:
        p = tf8.add_paragraph()
        p.text = t_item
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_SUCCESS
        p.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 9: Conclusion & Q&A
    # ----------------------------------------------------
    slide9 = add_dark_slide("Conclusion & Q&A")
    tb9 = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.8))
    tf9 = tb9.text_frame
    tf9.word_wrap = True

    p_end1 = tf9.paragraphs[0]
    p_end1.text = "THANK YOU!"
    p_end1.font.name = 'Arial'
    p_end1.font.size = Pt(46)
    p_end1.font.bold = True
    p_end1.font.color.rgb = ACCENT_CYAN
    p_end1.alignment = PP_ALIGN.CENTER

    p_end2 = tf9.add_paragraph()
    p_end2.text = "ARIA AI Voice Assistant is complete with Google OAuth, Email/Password Auth & SQL Database."
    p_end2.font.name = 'Arial'
    p_end2.font.size = Pt(18)
    p_end2.font.color.rgb = TEXT_LIGHT
    p_end2.alignment = PP_ALIGN.CENTER

    p_end3 = tf9.add_paragraph()
    p_end3.text = "\nQuestions & Answers (Q&A)"
    p_end3.font.name = 'Arial'
    p_end3.font.size = Pt(22)
    p_end3.font.bold = True
    p_end3.font.color.rgb = ACCENT_PURPLE
    p_end3.alignment = PP_ALIGN.CENTER

    prs.save("AI_Voice_Assistant_Presentation.pptx")
    print("[OK] Created updated AI_Voice_Assistant_Presentation.pptx successfully!")

if __name__ == '__main__':
    create_pptx_presentation()
